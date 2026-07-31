try:
    from ._imports import without_local_module_shadowing
except ImportError:
    from _imports import without_local_module_shadowing

with without_local_module_shadowing(__file__):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn

try:
    from ._layout import render_grid
    from .source import as_source
except ImportError:
    from _layout import render_grid
    from source import as_source


class PPTXExtractor(object):
    def __init__(self, filename, include_notes=True, include_charts=True, name=None):
        self.source = as_source(filename, name)
        self.filename = self.source.name
        self.include_notes = include_notes
        self.include_charts = include_charts
        self.slides = self._extract()
        self.text = self._render()
        self.provenance = self.source.provenance("pptx")

    def get_text(self):
        return self.text

    def get_structure(self):
        return {"type": "pptx", "provenance": self.provenance, "slides": self.slides}

    def get_slides(self):
        return self.slides

    def _extract(self):
        with self.source.open() as stream:
            presentation = Presentation(stream)
            slides = []
            for slide_index, slide in enumerate(presentation.slides, 1):
                self._slide_list_counters = {}
                blocks = self._shape_blocks(slide.shapes)
                blocks.sort(
                    key=lambda block: (
                        block.get("top", 0),
                        block.get("left", 0),
                        block.get("zOrder", 0),
                    )
                )
                notes = ""
                if self.include_notes and slide.has_notes_slide:
                    frame = slide.notes_slide.notes_text_frame
                    if frame is not None:
                        notes = frame.text.strip()
                slides.append(
                    {
                        "number": slide_index,
                        "name": slide.name,
                        "blocks": blocks,
                        "notes": notes,
                    }
                )
            return slides

    def _shape_blocks(self, shapes, parent_offset=(0, 0)):
        blocks = []
        for z_order, shape in enumerate(shapes):
            left = int(getattr(shape, "left", 0) or 0) + parent_offset[0]
            top = int(getattr(shape, "top", 0) or 0) + parent_offset[1]
            base = {
                "left": left,
                "top": top,
                "width": int(getattr(shape, "width", 0) or 0),
                "height": int(getattr(shape, "height", 0) or 0),
                "zOrder": z_order,
                "name": getattr(shape, "name", ""),
            }

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                blocks.extend(self._shape_blocks(shape.shapes, (left, top)))
                continue
            if getattr(shape, "has_table", False):
                block = dict(base)
                block.update(kind="table", table=self._table_data(shape.table))
                blocks.append(block)
                continue
            if self.include_charts and getattr(shape, "has_chart", False):
                block = dict(base)
                block.update(kind="chart", chart=self._chart_data(shape.chart))
                blocks.append(block)
                continue
            if getattr(shape, "has_text_frame", False):
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    if not paragraph.text:
                        continue
                    level = int(paragraph.level or 0)
                    paragraphs.append(
                        {
                            "text": paragraph.text,
                            "level": level,
                            "marker": self._pptx_list_marker(paragraph, level),
                        }
                    )
                if paragraphs:
                    block = dict(base)
                    block.update(kind="text", paragraphs=paragraphs)
                    blocks.append(block)
                continue

            # SmartArt and unsupported graphic frames can still carry a:t text.
            texts = [
                node.text
                for node in shape.element.iter(qn("a:t"))
                if node.text
            ]
            if texts:
                block = dict(base)
                block.update(kind="text", paragraphs=[{"text": value, "level": 0} for value in texts])
                blocks.append(block)
        return blocks

    def _pptx_list_marker(self, paragraph, level):
        properties = paragraph._p.pPr
        if properties is None:
            return ""
        bullet = properties.find(qn("a:buChar"))
        if bullet is not None:
            return (bullet.get("char") or "•") + " "
        auto_number = properties.find(qn("a:buAutoNum"))
        if auto_number is not None:
            key = level
            self._slide_list_counters[key] = self._slide_list_counters.get(key, 0) + 1
            for deeper in list(self._slide_list_counters):
                if deeper > level:
                    del self._slide_list_counters[deeper]
            return "{}. ".format(self._slide_list_counters[key])
        return ""

    @staticmethod
    def _table_data(table):
        rows = len(table.rows)
        cols = len(table.columns)
        cells = []
        for row_index in range(rows):
            for col_index in range(cols):
                cell = table.cell(row_index, col_index)
                if getattr(cell, "is_spanned", False):
                    continue
                cells.append(
                    {
                        "row": row_index,
                        "col": col_index,
                        "rowSpan": int(getattr(cell, "span_height", 1) or 1),
                        "colSpan": int(getattr(cell, "span_width", 1) or 1),
                        "text": cell.text.strip(),
                    }
                )
        return {
            "rows": rows,
            "cols": cols,
            "cells": cells,
            "columnWidths": [float(column.width) for column in table.columns],
        }

    @staticmethod
    def _chart_data(chart):
        title = ""
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text
        categories = []
        series = []
        try:
            if chart.plots and chart.plots[0].categories:
                categories = [str(category.label) for category in chart.plots[0].categories]
        except (AttributeError, TypeError, ValueError):
            pass
        for item in chart.series:
            try:
                values = [str(value) for value in item.values]
            except (AttributeError, TypeError, ValueError):
                values = []
            series.append({"name": str(item.name or ""), "values": values})
        return {"title": title, "categories": categories, "series": series}

    @staticmethod
    def _render_table(table):
        grid = [["" for _ in range(table["cols"])] for _ in range(table["rows"])]
        for cell in table["cells"]:
            value = cell["text"]
            if cell["rowSpan"] > 1 or cell["colSpan"] > 1:
                value += " (병합 {}×{})".format(cell["rowSpan"], cell["colSpan"])
            grid[cell["row"]][cell["col"]] = value
        return render_grid(
            grid,
            title="[표 {}×{}]".format(table["rows"], table["cols"]),
            column_weights=table["columnWidths"],
        )

    @staticmethod
    def _render_chart(chart):
        lines = ["[차트{}]".format(": " + chart["title"] if chart["title"] else "")]
        if chart["categories"]:
            lines.append("범주: " + ", ".join(chart["categories"]))
        for series in chart["series"]:
            lines.append(
                "{}: {}".format(series["name"] or "계열", ", ".join(series["values"]))
            )
        return "\n".join(lines)

    def _render(self):
        output = []
        for slide in self.slides:
            lines = ["[슬라이드 {}{}]".format(
                slide["number"], ": " + slide["name"] if slide["name"] else ""
            )]
            for block in slide["blocks"]:
                if block["kind"] == "table":
                    lines.append(self._render_table(block["table"]))
                elif block["kind"] == "chart":
                    lines.append(self._render_chart(block["chart"]))
                else:
                    for paragraph in block["paragraphs"]:
                        lines.append(
                            "  " * paragraph["level"]
                            + paragraph.get("marker", "")
                            + paragraph["text"]
                        )
            if slide["notes"]:
                lines.append("[발표자 노트]\n" + slide["notes"])
            output.append("\n\n".join(lines))
        return "\n\n".join(output)


def get_text(filename):
    print(PPTXExtractor(filename).get_text())


if __name__ == "__main__":
    import sys

    get_text(sys.argv[1])
